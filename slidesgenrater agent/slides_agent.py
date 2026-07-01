import json
import requests
import os
import re
import io
import copy
from duckduckgo_search import DDGS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
try:
    from lxml import etree
except ImportError:
    etree = None

import google.generativeai as genai
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
KIMI_API_KEY = os.getenv("KIMI_API_KEY")
if GEMINI_API_KEY:
    try: genai.configure(api_key=GEMINI_API_KEY)
    except: pass

THEMES = {
    "1": {"name": "Midnight Progressive (User Custom)", "bg": RGBColor(11, 18, 32), "primary": RGBColor(17, 24, 39), "accent": RGBColor(212, 175, 127), "text_dark": RGBColor(11, 18, 32), "text_light": RGBColor(249, 250, 251), "glass": RGBColor(15, 23, 42), "f_h": "Outfit", "f_b": "Outfit", "f_s": "Georgia"},
    "2": {"name": "Executive Black (Luxury)", "bg": RGBColor(10, 10, 10), "primary": RGBColor(20, 20, 20), "accent": RGBColor(255, 255, 255), "text_dark": RGBColor(0, 0, 0), "text_light": RGBColor(255, 255, 255), "glass": RGBColor(39, 39, 42), "f_h": "Outfit", "f_b": "Outfit", "f_s": "Georgia"},
    "3": {"name": "Premium White (Consulting)", "bg": RGBColor(255, 255, 255), "primary": RGBColor(245, 245, 245), "accent": RGBColor(37, 99, 235), "text_dark": RGBColor(17, 24, 39), "text_light": RGBColor(17, 24, 39), "glass": RGBColor(229, 231, 235), "f_h": "Outfit", "f_b": "Outfit", "f_s": "Georgia"},
    "4": {"name": "Navy Corporate (Serious Business)", "bg": RGBColor(15, 23, 42), "primary": RGBColor(30, 41, 59), "accent": RGBColor(56, 189, 248), "text_dark": RGBColor(15, 23, 42), "text_light": RGBColor(248, 250, 252), "glass": RGBColor(51, 65, 85), "f_h": "Outfit", "f_b": "Outfit", "f_s": "Georgia"},
    "5": {"name": "Luxury Purple (High-End)", "bg": RGBColor(26, 16, 44), "primary": RGBColor(42, 27, 77), "accent": RGBColor(245, 158, 11), "text_dark": RGBColor(26, 16, 44), "text_light": RGBColor(245, 243, 255), "glass": RGBColor(59, 47, 99), "f_h": "Outfit", "f_b": "Outfit", "f_s": "Georgia"},
    "11": {"name": "Bold Marketing (Power)", "bg": RGBColor(17, 24, 39), "primary": RGBColor(31, 41, 55), "accent": RGBColor(239, 68, 68), "text_dark": RGBColor(17, 24, 39), "text_light": RGBColor(255, 255, 255), "glass": RGBColor(55, 65, 81), "f_h": "Outfit", "f_b": "Outfit", "f_s": "Georgia"},
    "12": {"name": "Nature (Growth)", "bg": RGBColor(236, 253, 245), "primary": RGBColor(209, 250, 229), "accent": RGBColor(22, 163, 74), "text_dark": RGBColor(20, 83, 45), "text_light": RGBColor(20, 83, 45), "glass": RGBColor(187, 247, 208), "f_h": "Outfit", "f_b": "Outfit", "f_s": "Georgia"},
    "13": {"name": "Soft Education (Trust)", "bg": RGBColor(240, 249, 255), "primary": RGBColor(224, 242, 254), "accent": RGBColor(2, 132, 199), "text_dark": RGBColor(12, 74, 110), "text_light": RGBColor(12, 74, 110), "glass": RGBColor(186, 230, 253), "f_h": "Outfit", "f_b": "Outfit", "f_s": "Georgia"},
    "14": {"name": "Monochrome Pro (Ultra Clean)", "bg": RGBColor(255, 255, 255), "primary": RGBColor(250, 250, 250), "accent": RGBColor(0, 0, 0), "text_dark": RGBColor(0, 0, 0), "text_light": RGBColor(0, 0, 0), "glass": RGBColor(229, 229, 229), "f_h": "Outfit", "f_b": "Outfit", "f_s": "Georgia"},
    "15": {"name": "Cyberpunk Glow (Tech)", "bg": RGBColor(10, 10, 20), "primary": RGBColor(15, 15, 30), "accent": RGBColor(0, 255, 255), "text_dark": RGBColor(10, 10, 20), "text_light": RGBColor(200, 230, 255), "glass": RGBColor(20, 20, 50), "f_h": "Outfit", "f_b": "Outfit", "f_s": "Georgia"},
    "16": {"name": "Sunset Gold (Luxury)", "bg": RGBColor(28, 25, 23), "primary": RGBColor(41, 37, 36), "accent": RGBColor(234, 179, 8), "text_dark": RGBColor(28, 25, 23), "text_light": RGBColor(250, 250, 249), "glass": RGBColor(68, 64, 60), "f_h": "Outfit", "f_b": "Outfit", "f_s": "Georgia"}
}

# ==========================================
# DESIGN UTILITIES
# ==========================================
def set_background(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_background_decoration(slide, theme):
    """Adds ultra-premium procedural art to the background."""
    accent = theme["accent"]
    
    # 1. Radiant Geometric Wireframe (Institutional Tech vibe)
    for i in range(1, 15):
        # Radiating from bottom-left
        line = slide.shapes.add_connector(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.5), Inches(i * 1.1), Inches(0))
        line.line.color.rgb = accent; line.line.width = Pt(0.05)
        try: line.line.fill.transparency = 0.98
        except: pass
        
        # Radiating from top-right
        line2 = slide.shapes.add_connector(MSO_SHAPE.RECTANGLE, Inches(13.33), Inches(0), Inches(13.33 - i * 1.1), Inches(7.5))
        line2.line.color.rgb = accent; line2.line.width = Pt(0.05)
        try: line2.line.fill.transparency = 0.98
        except: pass

    # 2. Glassmorphism Gradient Overlay (Simulated)
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.5), Inches(13.33), Inches(1.0))
    overlay.fill.solid(); overlay.fill.fore_color.rgb = accent
    try: overlay.fill.transparency = 0.97
    except: pass
    overlay.line.width = 0

    # 3. Precision Corner Brackets
    # Top-Left
    add_shape(slide, accent, Inches(0.1), Inches(0.1), Inches(0.6), Inches(0.01), rounded=False, transparency=0.3, shadow=False)
    add_shape(slide, accent, Inches(0.1), Inches(0.1), Inches(0.01), Inches(0.6), rounded=False, transparency=0.3, shadow=False)
    # Bottom-Right
    add_shape(slide, accent, Inches(12.63), Inches(7.39), Inches(0.6), Inches(0.01), rounded=False, transparency=0.3, shadow=False)
    add_shape(slide, accent, Inches(13.23), Inches(6.8), Inches(0.01), Inches(0.6), rounded=False, transparency=0.3, shadow=False)

def add_shape(slide, color, left, top, width, height, transparency=0.0, line_color=None, shadow=True, rounded=True, border_width=1.0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    
    # Set fill
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if hasattr(shape.fill, 'transparency'):
        shape.fill.transparency = transparency
    
    # Premium Card Shadow Logic (Deep Floating Effect)
    if shadow:
        try:
            s = shape.shadow
            s.inherit = False
            s.visible = True
            s.blur_radius = Pt(28)
            s.distance = Pt(8)
            s.direction = 45
            s.transparency = 0.65
        except: pass
        
    # Modern Card Corner & Border Logic
    try:
        if rounded:
            shape.adjustments[0] = 0.06 # Precise rounding
            
        if line_color:
            shape.line.color.rgb = line_color
            shape.line.width = Pt(border_width)
            if transparency > 0:
                try: shape.line.fill.transparency = 0.4
                except: pass
        else:
            shape.line.width = 0
    except: pass
        
    return shape

def add_glass_box(slide, theme, accent, x, y, w, h):
    """Adds a sophisticated multi-layered glass box with a glowing border."""
    # 1. Subtle Outer Glow
    add_shape(slide, accent, Inches(x-0.03), Inches(y-0.03), Inches(w+0.06), Inches(h+0.06), transparency=0.9, rounded=True, shadow=True)
    # 2. Main Glass Body
    add_shape(slide, theme["glass"], Inches(x), Inches(y), Inches(w), Inches(h), transparency=0.25, rounded=True, line_color=accent, border_width=0.75)
    # 3. Precision Accent Bar
    add_shape(slide, accent, Inches(x + 0.3), Inches(y), Inches(1.5), Inches(0.04), transparency=0.1, rounded=False, shadow=False)

def add_textbox(slide, left, top, width, height, text, font_size=16, color=RGBColor(255,255,255), bold=False, align=PP_ALIGN.LEFT, font_name="Outfit", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]; p.alignment = align
    if text:
        run = p.add_run(); run.text = text; run.font.size = Pt(font_size)
        run.font.color.rgb = color; run.font.bold = bold; run.font.name = font_name
    return txBox


def adapt_font_size(text, base_size, min_size=12, max_size=72):
    """Dynamically adjust font size based on text length to prevent overflow or fill empty space."""
    length = len(str(text))
    # Aggressive scaling for large titles to prevent line wrapping into other elements
    if base_size > 30:
        if length > 100: size = base_size - 24
        elif length > 60: size = base_size - 16
        elif length > 35: size = base_size - 8
        elif length < 15: size = base_size + 8
        else: size = base_size
    else:
        if length > 300: size = base_size - 4
        elif length > 180: size = base_size - 2
        elif length < 35: size = base_size + 6
        elif length < 70: size = base_size + 4
        else: size = base_size
    return max(min_size, min(max_size, size))

# ==========================================
# RENDERERS (MULTI-LAYOUT SYSTEM)
# ==========================================
def render_hero_slide(slide, theme, item, img=None):
    """Ultra-Premium Geometric Title Slide."""
    set_background(slide, theme["bg"])
    accent = theme["accent"]
    add_background_decoration(slide, theme)
    
    # 1. Asymmetrical Geometric Brand Elements (Contained)
    add_shape(slide, accent, Inches(-0.5), Inches(0.5), Inches(7), Inches(7), transparency=0.96, rounded=True, shadow=False)
    add_shape(slide, accent, Inches(11.0), Inches(0.5), Inches(4), Inches(4), transparency=0.98, rounded=True, shadow=False)
    
    # 2. Main Title Glass Container
    add_glass_box(slide, theme, accent, 1.0, 2.0, 11.33, 3.5)
    
    # 3. Branding Marker
    add_shape(slide, accent, Inches(1.5), Inches(2.0), Inches(2.5), Inches(0.06), transparency=0.2, rounded=False)
    add_textbox(slide, 1.5, 1.5, 4.0, 0.4, "ADVISORY STRATEGIC REPORT", 10, accent, True, PP_ALIGN.LEFT, font_name=theme["f_h"])
    
    # 4. Content
    title_size = adapt_font_size(item["title"], 54, max_size=62)
    add_textbox(slide, 1.2, 2.5, 10.9, 1.8, item["title"].upper(), title_size, theme["text_light"], True, PP_ALIGN.CENTER, font_name=theme["f_h"], anchor=MSO_ANCHOR.MIDDLE)
    
    lead = item.get("lead_in", "Executive Strategy Formulation & Analysis")
    add_textbox(slide, 1.2, 4.4, 10.9, 0.6, lead, 22, accent, False, PP_ALIGN.CENTER, font_name=theme["f_s"])
    
    # 5. Footer Marker
    add_textbox(slide, 0.5, 7.0, 4.0, 0.3, "CONFIDENTIAL // AI SYNTHESIZED DATA", 8, accent, True, font_name=theme["f_h"])

def render_metrics_slide(slide, theme, item):
    set_background(slide, theme["bg"])
    accent = theme["accent"]
    txt_c = theme["text_light"]
    
    # Label Tag
    add_textbox(slide, 1.0, 0.2, 11.3, 0.4, "PERFORMANCE METRICS", 12, accent, True, PP_ALIGN.CENTER, font_name=theme["f_s"])

    # Title
    title_size = adapt_font_size(item["title"], 36)
    add_textbox(slide, 1.0, 0.5, 11.3, 1.0, item["title"].upper(), title_size, accent, True, PP_ALIGN.CENTER, font_name=theme["f_h"])
    
    metrics = item.get("metrics", [])[:4]
    if not metrics: # Fallback to bullets if metrics not provided
        render_density_slide(slide, theme, item)
        return

    count = len(metrics)
    col_w = (12.33 / count) - 0.4
    for i, m in enumerate(metrics):
        left = 0.5 + (i * (col_w + 0.4))
        # Ultra-Modern Metric Card
        add_glass_box(slide, theme, accent, left, 2.5, col_w, 3.5)
        
        # Value
        val = str(m.get("value", "0"))
        add_textbox(slide, left, 3.0, col_w, 1.5, val, 52, accent, True, PP_ALIGN.CENTER, font_name=theme["f_h"], anchor=MSO_ANCHOR.MIDDLE)
        
        # Label
        lab = str(m.get("label", "Metric"))
        add_textbox(slide, left + 0.1, 4.5, col_w - 0.2, 1.2, lab, 16, txt_c, True, PP_ALIGN.CENTER, font_name=theme["f_b"], anchor=MSO_ANCHOR.TOP)

def render_process_slide(slide, theme, item):
    set_background(slide, theme["bg"])
    accent = theme["accent"]
    txt_c = theme["text_light"]
    
    # Label Tag
    add_textbox(slide, 1.0, 0.2, 11.3, 0.4, "EXECUTION ROADMAP", 12, accent, True, PP_ALIGN.CENTER, font_name=theme["f_s"])

    title_size = adapt_font_size(item["title"], 36)
    add_textbox(slide, 1.0, 0.5, 11.3, 1.0, item["title"].upper(), title_size, accent, True, PP_ALIGN.CENTER, font_name=theme["f_h"])
    
    steps = item.get("steps", [])[:5]
    if not steps:
        render_density_slide(slide, theme, item)
        return

    count = len(steps)
    step_w = (12.33 / count) - 0.3
    for i, s in enumerate(steps):
        left = 0.5 + (i * (step_w + 0.3))
        
        # Step Circle/Indicator
        indicator = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + (step_w/2) - 0.4), Inches(2.2), Inches(0.8), Inches(0.8))
        indicator.fill.solid(); indicator.fill.fore_color.rgb = accent
        indicator.line.width = 0
        
        # Step Number
        add_textbox(slide, left + (step_w/2) - 0.4, 2.2, 0.8, 0.8, str(i+1), 24, theme["text_dark"], True, PP_ALIGN.CENTER, font_name=theme["f_h"], anchor=MSO_ANCHOR.MIDDLE)
        
        # Connector Line
        if i < count - 1:
            line = slide.shapes.add_connector(MSO_SHAPE.RECTANGLE, Inches(left + (step_w/2) + 0.4), Inches(2.6), Inches(left + step_w + 0.3 + (step_w/2) - 0.4), Inches(2.6))
            line.line.color.rgb = accent; line.line.width = Pt(2)
            
        # Content Card
        add_shape(slide, theme["glass"], Inches(left), Inches(3.5), Inches(step_w), Inches(3.0), line_color=accent, shadow=True)
        add_textbox(slide, left + 0.1, 3.6, step_w - 0.2, 2.8, str(s), 16, txt_c, False, PP_ALIGN.CENTER, font_name=theme["f_b"], anchor=MSO_ANCHOR.MIDDLE)

def render_comparison_slide(slide, theme, item):
    set_background(slide, theme["bg"])
    accent = theme["accent"]
    txt_c = theme["text_light"]
    
    # Label Tag
    add_textbox(slide, 1.0, 0.2, 11.3, 0.4, "STRATEGIC COMPARISON", 12, accent, True, PP_ALIGN.CENTER, font_name=theme["f_s"])

    title_size = adapt_font_size(item["title"], 36)
    add_textbox(slide, 1.0, 0.5, 11.3, 1.0, item["title"].upper(), title_size, accent, True, PP_ALIGN.CENTER, font_name=theme["f_h"])
    
    comp = item.get("comparison", {})
    left_data = comp.get("left", {"title": "CURRENT", "bullets": ["No data"]})
    right_data = comp.get("right", {"title": "PROPOSED", "bullets": ["No data"]})
    
    # Left Card
    add_shape(slide, theme["primary"], Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5), line_color=accent, shadow=True)
    add_textbox(slide, 1.0, 2.2, 5.1, 0.8, left_data["title"].upper(), 24, accent, True, PP_ALIGN.CENTER, font_name=theme["f_h"])
    l_text = "\n\n".join(["• " + str(b) for b in left_data.get("bullets", [])[:4]])
    add_textbox(slide, 1.0, 3.0, 5.1, 3.3, l_text, 16, txt_c, False, PP_ALIGN.LEFT, font_name=theme["f_b"])
    
    # Right Card
    add_shape(slide, theme["glass"], Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5), line_color=accent, shadow=True)
    add_textbox(slide, 7.2, 2.2, 5.1, 0.8, right_data["title"].upper(), 24, accent, True, PP_ALIGN.CENTER, font_name=theme["f_h"])
    r_text = "\n\n".join(["• " + str(b) for b in right_data.get("bullets", [])[:4]])
    add_textbox(slide, 7.2, 3.0, 5.1, 3.3, r_text, 16, txt_c, False, PP_ALIGN.LEFT, font_name=theme["f_b"])
    
    # VS Circle
    vs = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.3), Inches(3.9), Inches(0.8), Inches(0.8))
    vs.fill.solid(); vs.fill.fore_color.rgb = accent
    add_textbox(slide, 6.3, 3.9, 0.8, 0.8, "VS", 20, theme["text_dark"], True, PP_ALIGN.CENTER, font_name=theme["f_h"], anchor=MSO_ANCHOR.MIDDLE)

def render_quote_slide(slide, theme, item):
    set_background(slide, theme["bg"])
    accent = theme["accent"]
    txt_c = theme["text_light"]
    
    # Large Quote Marks
    add_textbox(slide, 1.0, 1.5, 2.0, 2.0, "“", 120, accent, True, PP_ALIGN.LEFT, font_name=theme["f_s"])
    
    quote_text = item.get("quote", item.get("lead_in", "Excellence is not an act, but a habit."))
    q_size = adapt_font_size(quote_text, 36, max_size=48)
    add_textbox(slide, 2.0, 2.5, 9.33, 3.0, quote_text, q_size, txt_c, True, PP_ALIGN.CENTER, font_name=theme["f_s"], anchor=MSO_ANCHOR.MIDDLE)
    
    author = item.get("author", "Strategic Insight")
    add_textbox(slide, 2.0, 5.5, 9.33, 1.0, f"— {author}", 24, accent, False, PP_ALIGN.RIGHT, font_name=theme["f_h"])

def render_density_slide(slide, theme, item, img=None):
    set_background(slide, theme["bg"])
    add_background_decoration(slide, theme)
    txt_c = theme["text_light"]
    accent = theme["accent"]
    
    # Top Label Tag (Serif)
    add_textbox(slide, 0.6, 0.1, 5.0, 0.4, "STRATEGIC OVERVIEW", 12, accent, True, PP_ALIGN.LEFT, font_name=theme["f_s"])

    # Header with dynamic sizing
    title_size = adapt_font_size(item["title"], 34, max_size=38)
    add_textbox(slide, 0.6, 0.4, 11.8, 0.8, item["title"].upper(), title_size, accent, True, PP_ALIGN.LEFT, font_name=theme["f_h"])
    
    # Title Branding Line
    add_shape(slide, accent, Inches(0.6), Inches(1.15), Inches(4.0), Inches(0.03), rounded=False, transparency=0.2, shadow=False)
    
    # Image placement if available
    image_active = False
    if img:
        try:
            # Card frame for image
            add_shape(slide, theme["primary"], Inches(8.5), Inches(1.3), Inches(4.3), Inches(5.7), shadow=True, line_color=accent)
            slide.shapes.add_picture(img, Inches(8.7), Inches(1.5), width=Inches(3.9))
            content_w = 7.5
        except: content_w = 12.1
    else:
        content_w = 12.1

    # Lead-in text box (Executive Insight) - DYNAMIC & OVERFLOW-SAFE
    lead = item.get("lead_in", "")
    if lead:
        # Generous height calculation to prevent overflow
        txt_len = len(str(lead))
        if txt_len < 100: box_h = 1.0; l_size = 26
        elif txt_len < 200: box_h = 1.3; l_size = 22
        else: box_h = 1.7; l_size = 19
        
        # Premium Glow-Bar & Glass Container
        add_glass_box(slide, theme, accent, 0.6, 1.3, content_w, box_h)
        # Accent bar on the side
        add_shape(slide, accent, Inches(0.6), Inches(1.3), Inches(0.04), Inches(box_h), rounded=False, transparency=0.1)
        
        # Subtitle (Small Caps / Serif)
        add_textbox(slide, 0.8, 1.35, 4.0, 0.3, "EXECUTIVE SUMMARY", 10, accent, True, PP_ALIGN.LEFT, font_name=theme["f_h"])
        # Content (Elegant Serif) - Centered vertically with buffer
        add_textbox(slide, 0.8, 1.6, content_w - 0.6, box_h - 0.4, lead, l_size, txt_c, False, PP_ALIGN.LEFT, font_name=theme["f_s"], anchor=MSO_ANCHOR.MIDDLE)
        
        top_pos = 1.3 + box_h + 0.2
    else:
        top_pos = 1.3
    
    bullets = item.get("bullets", [])
    
    for i, b in enumerate(bullets[:6]): # Limit to 6
        y = top_pos
        
        # Content-Aware height for bullets (Overflow-Safe)
        b_txt = str(b)
        b_len = len(b_txt)
        b_h = 0.9 if b_len < 150 else 1.2
        
        # STRICT VERTICAL LIMIT: Prevent going into the bottom accent bar (starts at 6.5)
        if y + b_h > 6.4:
            break
            
        bg = theme["primary"] if i % 2 == 0 else theme["glass"]
        add_glass_box(slide, theme, accent, 0.6, y, content_w, b_h)
        
        # Bullet Text - Increased for readability
        b_size = adapt_font_size(b_txt, 20, max_size=24)
        add_textbox(slide, 0.9, y, content_w - 0.6, b_h, "• " + b_txt, b_size, txt_c, False, PP_ALIGN.LEFT, font_name=theme["f_b"], anchor=MSO_ANCHOR.MIDDLE)
        
        # Shift next bullet down based on current height
        top_pos += (b_h + 0.1)

# ==========================================
# ULTRA-PIPELINE (UNIFIED)
# ==========================================
def clean_vc_text(text):
    if not text: return ""
    lines = text.split("\n"); cleaned = []
    forbidden = [r'^Page\s+\d+$', r'Table of Contents', r'Update Field']
    for line in lines:
        s = line.strip()
        if any(re.match(f, s, re.I) for f in forbidden): continue 
        if len(s) < 3: continue
        cleaned.append(s)
    return "\n".join(cleaned)

def generate_slide_content(topic, num_slides=10, detail_level="detailed", subtopics="", subtopics_behavior="only", document_text="", task_updater=None):
    text = clean_vc_text(document_text)
    if task_updater: task_updater("Initiating Advanced Content Synthesis...")
    
    # Check for Kimi First if requested
    if KIMI_API_KEY:
        try:
            if task_updater: task_updater("Structuring core intelligence...")
            content = call_kimi_api(topic, num_slides, detail_level, subtopics, subtopics_behavior, text)
            if content:
                blueprint = safe_parse(content)
                if blueprint: return validate_slide_content(blueprint, text, topic, num_slides)
        except Exception as e:
            print(f"[KIMI] Fallback to Gemini: {e}")

    # Priority model list including user preference
    env_model = os.getenv("MODEL_NAME")
    models = [env_model] if env_model else []
    # Use standard names that are most likely to work across versions
    models.extend(["gemini-1.5-flash", "gemini-1.5-pro", "gemini-1.0-pro"])
    
    for m in models:
        if not m: continue
        try:
            if task_updater: task_updater("Drafting high-fidelity slides...")
            # Explicitly try to use the model name without 'models/' prefix if it's there
            model_name = m.replace("models/", "")
            model = genai.GenerativeModel(model_name)
            
            # Content style and behavior logic
            if detail_level == "detailed":
                content_style = "Insight-dense, analytical, and verbose. Each bullet point MUST be a full, complex paragraph (40-60 words each) packed with data and strategic depth. Each lead_in MUST be a substantial executive summary."
            else:
                content_style = "Short, high-impact, and punchy bullet points."
            subtopic_instruction = f"Strictly focus on these pillars: {subtopics}." if subtopics and subtopics_behavior == "only" else f"Ensure inclusion of: {subtopics}." if subtopics else ""

            prompt = (
                f"Topic: {topic}\n"
                f"Context: {text[:20000]}\n"
                f"Requirements: {num_slides} distinct slides. Each slide MUST cover a unique sub-aspect of the topic with zero content overlap. DO NOT repeat bullet points or themes across different slides. {subtopic_instruction}\n"
                f"Content Depth: {content_style}.\n"
                f"Variety is key. Use different layouts where appropriate. Ensure each slide adds new value."
            )
            
            # Use JSON mode for 1.5 models if possible
            gen_config = {"response_mime_type": "application/json"} if "1.5" in m else {}
            
            resp = model.generate_content(prompt, generation_config=gen_config)
            blueprint = safe_parse(resp)
            
            if blueprint and len(blueprint) >= 1:
                if task_updater: task_updater("Finalizing content architecture...")
                return validate_slide_content(blueprint, text, topic, num_slides)
        except Exception as e:
            if task_updater: task_updater("Optimizing generation sequence...")
            print(f"[AI ERROR] {m}: {e}")
            continue

    if task_updater: task_updater("Synthesizing strategic context from source...")
    return get_investor_fallback(topic, num_slides, text, detail_level)

def fetch_image(query):
    """Fetches a professional image URL related to the query."""
    try:
        with DDGS() as ddgs:
            results = list(ddgs.images(query, max_results=1))
            if results:
                return results[0]['image']
    except:
        pass
    return None

def validate_slide_content(slides, text, topic, num_slides):
    """Clean and structure AI-generated slides without losing their content."""
    validated = []
    for i in range(num_slides):
        if i < len(slides):
            s = copy.deepcopy(slides[i])
        else:
            # Dynamic filler pool based on standard business pillars to prevent repetition
            fillers = [
                {
                    "title": f"Strategic Expansion: {topic}",
                    "lead_in": f"Scaling the core technical and operational foundations of {topic} to achieve global market dominance.",
                    "bullets": [
                        f"Scalability: Identifying key growth levers for {topic} in emerging markets.",
                        f"Competitive Moat: Establishing a sustainable advantage through technical innovation.",
                        f"Operational Excellence: Optimizing the core delivery pipeline for maximum efficiency.",
                        f"Risk Mitigation: Proactively addressing regulatory and market volatility."
                    ]
                },
                {
                    "title": f"Future Outlook & Vision",
                    "lead_in": f"Anticipating the long-term evolution of {topic} and positioning for sustained leadership.",
                    "bullets": [
                        f"Horizon 2 Strategy: Exploring adjacent markets and integration opportunities.",
                        f"R&D Pipeline: Prioritizing high-impact innovations that redefine the {topic} landscape.",
                        f"Sustainability: Integrating environmental and social governance into the core business model.",
                        f"Talent Acquisition: Building an elite global team to drive the next phase of growth."
                    ]
                },
                {
                    "title": f"Financial Resilience & ROI",
                    "lead_in": f"Optimizing capital allocation and revenue streams to ensure long-term stability and shareholder value.",
                    "bullets": [
                        f"Revenue Diversification: Identifying new monetization channels within the {topic} ecosystem.",
                        f"Margin Expansion: Leveraging automation and scale to drive operational efficiencies.",
                        f"Capital Efficiency: Rigorous assessment of investment returns across all business units.",
                        f"Risk-Adjusted Growth: Balancing aggressive expansion with prudent financial management."
                    ]
                }
            ]
            filler_idx = (i - len(slides)) % len(fillers)
            s = copy.deepcopy(fillers[filler_idx])
        
        # Ensure layout field exists
        if "layout" not in s: s["layout"] = "standard"
        
        # Ensure fields exist and use existing content where possible
        if "title" not in s or not s["title"]: s["title"] = f"Strategic Depth: {topic}"
        if "lead_in" not in s or not s["lead_in"]: 
            s["lead_in"] = f"An analytical deep-dive into the core pillars that define the {topic} ecosystem's maturity."
        
        if i == 0:
            s["type"] = "hero"
            s["title"] = str(topic).upper()
        else:
            s["type"] = s["layout"]
            
        validated.append(s)
    return validated

def get_investor_fallback(topic, num_slides, document_text="", detail_level="standard"):
    """A DYNAMIC fallback that generates MUCH longer text for 'detailed' mode."""
    slides = []
    import random
    
    # 1. Extract and Clean Snippets
    raw_snippets = [s.strip() for s in document_text.replace('\n', '.').split('.') if len(s.strip()) > 35]
    doc_snippets = []
    seen_snips = set()
    for s in raw_snippets:
        key = s[:50].lower()
        if key not in seen_snips:
            doc_snippets.append(s)
            seen_snips.add(key)

    if len(doc_snippets) < 10:
        doc_snippets.extend([
            f"Optimizing {topic} through advanced strategic alignment and operational excellence.",
            f"The {topic} roadmap focuses on scalable infrastructure and long-term sustainability.",
            f"Enhancing the core value proposition of {topic} for all primary stakeholders.",
            f"Market dynamics for {topic} require an agile and data-driven approach to dominance.",
            f"Phase-based execution of {topic} initiatives ensures consistent performance and ROI.",
            f"Analyzing the competitive landscape reveals key entry points for {topic} disruption.",
            f"Strategic resource allocation is critical for the global expansion of {topic}."
        ])
    
    # Hero Slide
    slides.append({
        "title": topic.upper(),
        "lead_in": doc_snippets[0][:250],
        "bullets": ["Executive Roadmap", "Strategic Positioning", "Financial Outlook", "Operational Excellence"]
    })
    
    # 2. Detailed Mode Logic: Combine snippets for length
    import random
    def get_varied_bullets(count, start_offset):
        bullets = []
        for j in range(count):
            idx = (start_offset + j * 3) % len(doc_snippets)
            s1 = doc_snippets[idx]
            if detail_level == "detailed":
                # Merge two snippets for 'longer' mode
                idx2 = (idx + 7) % len(doc_snippets) # Higher offset for variety
                s2 = doc_snippets[idx2]
                bullets.append(f"{s1} Furthermore, {s2[0].lower() + s2[1:]}")
            else:
                bullets.append(s1[:160])
        return bullets

    # 3. Randomized Slide Generation for variety
    layout_pool = ["standard", "metrics", "process", "comparison", "quote"]
    random.shuffle(layout_pool)
    
    for i in range(1, num_slides):
        base_offset = (i * 19 + len(topic)) % len(doc_snippets)
        layout = layout_pool[i % len(layout_pool)]
        
        if layout == "standard":
            slides.append({
                "title": f"STRATEGIC LANDSCAPE (PHASE {i})",
                "layout": "standard",
                "lead_in": doc_snippets[base_offset][:200],
                "bullets": get_varied_bullets(4, base_offset + 1)
            })
        elif layout == "process":
            slides.append({
                "title": f"EXECUTION FRAMEWORK: {topic.upper()}",
                "layout": "process",
                "lead_in": "Key operational milestones extracted from document analysis.",
                "steps": [s[:130] for s in get_varied_bullets(4, base_offset + 2)]
            })
        elif layout == "metrics":
            slides.append({
                "title": "CORE PERFORMANCE METRICS",
                "layout": "metrics",
                "lead_in": "Analytical targets based on current strategic roadmap.",
                "metrics": [
                    {"value": f"{random.randint(70,99)}%", "label": doc_snippets[(base_offset + 1) % len(doc_snippets)][:45]},
                    {"value": f"Q{random.randint(1,4)}", "label": doc_snippets[(base_offset + 3) % len(doc_snippets)][:45]}
                ]
            })
        elif layout == "comparison":
            slides.append({
                "title": "STRATEGIC POSITIONING COMPARISON",
                "layout": "comparison",
                "comparison": {
                    "left": {"title": "TRADITIONAL", "bullets": get_varied_bullets(2, base_offset + 3)},
                    "right": {"title": "OPTIMIZED", "bullets": get_varied_bullets(2, base_offset + 5)}
                }
            })
        else: # Quote
            slides.append({
                "title": "GUIDING VISION",
                "layout": "quote",
                "quote": doc_snippets[base_offset][:150],
                "author": "Strategic Analysis"
            })
            
    return validate_slide_content(slides, document_text, topic, num_slides)
            
    return validate_slide_content(slides, document_text, topic, num_slides)

def create_presentation(slides_data, topic, theme, detail_level, include_images=False, output_file="presentation.pptx", pdf_images=None, task_updater=None):
    if not slides_data: return False
    prs = Presentation(); prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    for i, item in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if task_updater: task_updater(f"Rendering Slide {i+1}...")
        
        img_path = None
        if pdf_images and i < len(pdf_images):
            img_path = pdf_images[i]
        elif include_images:
            if task_updater: task_updater(f"Finding visual for slide {i+1}...")
            img_url = fetch_image(f"{topic} {item['title']}")
            if img_url:
                try:
                    r = requests.get(img_url, stream=True, timeout=5)
                    if r.status_code == 200:
                        img_path = io.BytesIO(r.content)
                except: pass

        st = item.get("type", "standard").lower()
        if st == "hero": render_hero_slide(slide, theme, item, img_path)
        elif st == "metrics": render_metrics_slide(slide, theme, item)
        elif st == "process": render_process_slide(slide, theme, item)
        elif st == "comparison": render_comparison_slide(slide, theme, item)
        elif st == "quote": render_quote_slide(slide, theme, item)
        else: render_density_slide(slide, theme, item, img_path)
        
    try:
        prs.save(output_file)
        return True
    except Exception as e:
        print(f"[ERROR] Save failed: {e}")
        return False

def call_kimi_api(topic, num_slides, detail_level, subtopics, subtopics_behavior, context):
    """Moonshot AI (Kimi) implementation for high-quality content."""
    model_name = os.getenv("MODEL_NAME", "moonshot-v1-8k")
    
    if detail_level == "detailed":
        content_style = "Insight-dense, analytical, and verbose. Each bullet point MUST be a full, complex paragraph (40-60 words each) packed with data and strategic depth."
    else:
        content_style = "Short, high-impact, and punchy bullet points."
        
    subtopic_instruction = f"Strictly focus on these pillars: {subtopics}." if subtopics and subtopics_behavior == "only" else f"Ensure inclusion of: {subtopics}." if subtopics else ""

    prompt = (
        f"Topic: {topic}\nContext: {context[:15000]}\nRequirements: {num_slides} distinct slides.\n"
        f"Style: {content_style} {subtopic_instruction}\n"
        f"Return a JSON array where each object has: title, layout, lead_in, bullets, metrics, steps, comparison, quote, author."
    )
    
    headers = {"Authorization": f"Bearer {KIMI_API_KEY}", "Content-Type": "application/json"}
    payload = {
        "model": model_name,
        "messages": [
            {"role": "system", "content": "You are an elite presentation architect. Always return valid JSON only."},
            {"role": "user", "content": prompt}
        ],
        "response_format": {"type": "json_object"}
    }
    
    try:
        r = requests.post("https://api.moonshot.cn/v1/chat/completions", headers=headers, json=payload, timeout=45)
        if r.status_code == 200:
            return r.json()['choices'][0]['message']['content']
    except: pass
    return None

def safe_parse(resp):
    try:
        t = resp.text if hasattr(resp, 'text') else str(resp)
        cl = re.sub(r'```json\n?|\n?```', '', t).strip()
        m = re.search(r'\[\s*\{.*\}\s*\]', cl, re.DOTALL)
        return json.loads(m.group(0)) if m else json.loads(cl)
    except: return None
