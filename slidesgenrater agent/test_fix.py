"""Test script to verify the index out of range fix"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from slides_agent import render_hero_slide, render_standard_slide, render_kpi_slide, THEMES

def test_hero_slide():
    """Test hero slide rendering with multiple bullets"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme = THEMES["2"]
    
    item = {
        "title": "Test Presentation",
        "bullets": [
            "Bullet 1: Test content",
            "Bullet 2: More content",
            "Bullet 3: Additional info",
            "Bullet 4: Extra details",
            "Bullet 5: Final point",
            "Bullet 6: Last item"
        ],
        "type": "hero"
    }
    
    try:
        render_hero_slide(slide, theme, item)
        print("[PASS] Hero slide rendered successfully")
        return True
    except Exception as e:
        print(f"[FAIL] Hero slide failed: {e}")
        return False

def test_standard_slide():
    """Test standard slide rendering with multiple bullets"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme = THEMES["2"]
    
    item = {
        "title": "Test Slide",
        "bullets": [
            "Point 1: Detailed explanation",
            "Point 2: Another point",
            "Point 3: More info",
            "Point 4: Additional data",
            "Point 5: Final thoughts"
        ],
        "type": "standard"
    }
    
    try:
        render_standard_slide(slide, theme, item, is_even=True)
        print("[PASS] Standard slide rendered successfully")
        return True
    except Exception as e:
        print(f"[FAIL] Standard slide failed: {e}")
        return False

def test_kpi_slide():
    """Test KPI slide rendering"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme = THEMES["2"]
    
    item = {
        "title": "Key Metrics",
        "bullets": [
            "Revenue: $1.5M",
            "Growth: 45%",
            "Users: 10,000+"
        ],
        "type": "kpi"
    }
    
    try:
        render_kpi_slide(slide, theme, item)
        print("[PASS] KPI slide rendered successfully")
        return True
    except Exception as e:
        print(f"[FAIL] KPI slide failed: {e}")
        return False

if __name__ == "__main__":
    print("Testing slide rendering functions...")
    print("-" * 50)
    
    results = []
    results.append(test_hero_slide())
    results.append(test_standard_slide())
    results.append(test_kpi_slide())
    
    print("-" * 50)
    if all(results):
        print("[SUCCESS] All tests passed! The index out of range error should be fixed.")
    else:
        print("[FAILURE] Some tests failed")
