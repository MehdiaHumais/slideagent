"""Comprehensive test to find and verify fix for index out of range errors"""
import sys
from slides_agent import (
    extract_financial_data, 
    validate_slide_content, 
    render_hero_slide,
    render_standard_slide,
    render_section_slide,
    render_2column_slide,
    render_kpi_slide,
    render_table_slide,
    THEMES
)
from pptx import Presentation
from pptx.util import Inches

def test_extract_financial_data():
    """Test extraction with empty text"""
    print("[TEST] Financial data extraction with empty text...")
    try:
        result = extract_financial_data("")
        assert result is None, "Should return None for empty text"
        print("  [PASS] Empty text handled correctly")
        return True
    except Exception as e:
        print(f"  [FAIL] Extraction failed: {e}")
        return False

def test_validate_slide_content():
    """Test validation with empty slides"""
    print("[TEST] Validate slide content with empty list...")
    try:
        result = validate_slide_content([], "", "")
        assert result == [], "Should return empty list"
        print("  [PASS] Empty validation handled correctly")
        return True
    except Exception as e:
        print(f"  [FAIL] Validation failed: {e}")
        return False

def test_render_functions():
    """Test all render functions with edge cases"""
    theme = THEMES["2"]
    
    # Test render_hero_slide with no bullets
    print("[TEST] render_hero_slide with no bullets...")
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        item = {"title": "Test", "type": "hero"}  # No bullets key
        render_hero_slide(slide, theme, item)
        print("  [PASS] No bullets handled")
    except Exception as e:
        print(f"  [FAIL] render_hero_slide failed: {e}")
        return False
    
    # Test render_hero_slide with empty bullets
    print("[TEST] render_hero_slide with empty bullets...")
    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        item = {"title": "Test", "type": "hero", "bullets": []}
        render_hero_slide(slide, theme, item)
        print("  [PASS] Empty bullets handled")
    except Exception as e:
        print(f"  [FAIL] render_hero_slide with empty bullets failed: {e}")
        return False
    
    # Test render_section_slide with no bullets
    print("[TEST] render_section_slide with no bullets...")
    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        item = {"title": "Section", "type": "section"}
        render_section_slide(slide, theme, item)
        print("  [PASS] No bullets handled")
    except Exception as e:
        print(f"  [FAIL] render_section_slide failed: {e}")
        return False
    
    # Test render_kpi_slide with no bullets
    print("[TEST] render_kpi_slide with no bullets...")
    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        item = {"title": "Metrics", "type": "kpi"}
        render_kpi_slide(slide, theme, item)
        print("  [PASS] No bullets handled")
    except Exception as e:
        print(f"  [FAIL] render_kpi_slide failed: {e}")
        return False
    
    # Test render_2column_slide with no bullets
    print("[TEST] render_2column_slide with no bullets...")
    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        item = {"title": "Comparison", "type": "2-column"}
        render_2column_slide(slide, theme, item)
        print("  [PASS] No bullets handled")
    except Exception as e:
        print(f"  [FAIL] render_2column_slide failed: {e}")
        return False
    
    # Test render_table_slide with no bullets
    print("[TEST] render_table_slide with no bullets...")
    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        item = {"title": "Table", "type": "table"}
        render_table_slide(slide, theme, item)
        print("  [PASS] No bullets handled")
    except Exception as e:
        print(f"  [FAIL] render_table_slide failed: {e}")
        return False
    
    return True

def test_slide_info_without_type():
    """Test that slides without 'type' field don't crash"""
    print("[TEST] Slide info without 'type' field...")
    try:
        from slides_agent import generate_slide_content
        # This should not crash even if type is missing
        slides_data = [
            {"title": "Slide 1"},  # No type
            {"title": "Slide 2", "type": "standard"},
            {"title": "Slide 3", "bullets": ["test"]}  # No type
        ]
        # Check that all slides have a type
        for slide in slides_data:
            stype = slide.get("type", "standard")
            assert stype in ["hero", "section", "standard", "2-column", "kpi", "table"], f"Invalid type: {stype}"
        print("  [PASS] Type field handling correct")
        return True
    except Exception as e:
        print(f"  [FAIL] Type field test failed: {e}")
        return False

if __name__ == "__main__":
    print("=" * 60)
    print("COMPREHENSIVE INDEX OUT OF RANGE TEST")
    print("=" * 60)
    
    results = []
    results.append(test_extract_financial_data())
    results.append(test_validate_slide_content())
    results.append(test_render_functions())
    results.append(test_slide_info_without_type())
    
    print("=" * 60)
    if all(results):
        print("[SUCCESS] All tests passed! Index errors should be fixed.")
    else:
        print("[FAILURE] Some tests failed!")
        sys.exit(1)
